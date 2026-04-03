"""Build reference.pptx from eBots deck: structural slides + S4 pattern slides.

Slide order:
  1. S1 Title
  2. S2 Executive Summary
  3. S3 Divider
  4. S4 base (empty content slide)
  5-12. S4 patterns (standard-content, two-column, highlights, data-table, sidebar, testimonials, timeline, matrix)
  13. S5 Investment Criteria
  14. S6 Appendix
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# === Step 1: Build base structural slides from eBots deck ===
prs = Presentation(".claude/skills/investment-memo/assets/eBots_presentation.pptx")

ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
r_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

# Keep S1(0), S2(1), S3(2), S4(3), S21(20), S22(21)
# S21 and S22 will end up after patterns due to slide ordering
keep = {0, 1, 2, 3, 20, 21}
sldIdLst = prs._element.find(f"{ns}sldIdLst")
slide_ids = list(sldIdLst)
for i in range(len(slide_ids) - 1, -1, -1):
    if i not in keep:
        sldId = slide_ids[i]
        rId = sldId.get(r_ns + "id")
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

def replace_run_text(shape, new_text):
    if shape.has_text_frame and shape.text_frame.paragraphs:
        para = shape.text_frame.paragraphs[0]
        if para.runs:
            para.runs[0].text = new_text
            for run in para.runs[1:]:
                run._r.getparent().remove(run._r)
        for p in shape.text_frame.paragraphs[1:]:
            p._p.getparent().remove(p._p)

# S1: Title
slide = prs.slides[0]
for sh in slide.shapes:
    try:
        if sh.placeholder_format.idx == 0:
            replace_run_text(sh, "Investment Memo: Enter agenda nickname here")
    except (ValueError, AttributeError):
        if sh.has_text_frame and sh.shape_type == 17:
            replace_run_text(sh, "Month Year")
print("S1 done")

# S2: Exec Summary
slide = prs.slides[1]
for sh in slide.shapes:
    if sh.name == "Rectangle 4":
        replace_run_text(sh, "Enter agenda nickname here")
    elif sh.name == "Rectangle 5":
        if sh.has_text_frame and sh.text_frame.paragraphs:
            first_para = sh.text_frame.paragraphs[0]
            if first_para.runs:
                first_para.runs[0].text = "Enter company background and supporting facts here. Must include Investment Opportunity section."
                for run in first_para.runs[1:]:
                    run._r.getparent().remove(run._r)
            for p in sh.text_frame.paragraphs[1:]:
                p._p.getparent().remove(p._p)
print("S2 done")

# S3: Divider
slide = prs.slides[2]
for sh in slide.shapes:
    if sh.name == "TextBox 5":
        replace_run_text(sh, "Enter section name here")
print("S3 done")

# S4: Content base — strip content, keep chrome
slide = prs.slides[3]
shapes_to_remove = []
for sh in slide.shapes:
    is_ph = False
    try:
        sh.placeholder_format.idx
        is_ph = True
    except (ValueError, AttributeError):
        pass
    keep_shape = is_ph or sh.name == "think-cell data - do not delete" or sh.name == "Rectangle: Rounded Corners 43"
    if not keep_shape:
        shapes_to_remove.append(sh)
for sh in shapes_to_remove:
    sh._element.getparent().remove(sh._element)
for sh in slide.shapes:
    try:
        if sh.placeholder_format.idx == 0:
            replace_run_text(sh, "Sentence-form slide title goes here")
    except (ValueError, AttributeError):
        if sh.name == "Rectangle: Rounded Corners 43":
            replace_run_text(sh, "Enter section tag")
print(f"S4 base done ({len(slide.shapes)} shapes)")

# === Step 2: Add S4 pattern slides ===
content_layout = None
for layout in prs.slide_layouts:
    if layout.name == "Content":
        content_layout = layout
        break

def add_tag(slide, text="Tag"):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.66), Inches(0.05), Inches(2.51), Inches(0.28)
    )
    p = shape.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.bold = True

def add_source(slide, text="Source: /1 Source description here"):
    box = slide.shapes.add_textbox(Inches(0.47), Inches(6.63), Inches(12.3), Inches(0.5))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(9)

def set_title(slide, text):
    for sh in slide.shapes:
        try:
            if sh.placeholder_format.idx == 0:
                sh.text = text
                return
        except (ValueError, AttributeError):
            pass

def add_bold_bullets(slide, bullets, left=0.67, top=1.6, width=11.8, height=4.8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (bold, regular) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = bold
        r1.font.size = Pt(14)
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = regular
        r2.font.size = Pt(14)
        p.space_after = Pt(8)
    return box

def add_callout_box(slide, text, left, top, width, height):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x39, 0x33, 0x35)
    return box


# === 1. standard-content ===
slide = prs.slides.add_slide(content_layout)
set_title(slide, "Sentence-form title stating the key takeaway of this slide")
add_tag(slide, "Tag")
add_bold_bullets(slide, [
    ("Bold claim phrase here", " supporting detail and evidence with specific data points"),
    ("Second key point", " with metric or proof point that strengthens the argument"),
    ("Third supporting argument", " additional evidence drawn from the memo with sourced numbers"),
    ("Fourth data point if needed", " keeps the slide complete without overcrowding"),
])
add_source(slide)
print("standard-content done")


# === 2. two-column ===
slide = prs.slides.add_slide(content_layout)
set_title(slide, "Sentence-form title comparing two parallel concepts")
add_tag(slide, "Tag")

# Left header
lh = slide.shapes.add_textbox(Inches(0.67), Inches(1.6), Inches(5.5), Inches(0.4))
r = lh.text_frame.paragraphs[0].add_run()
r.text = "LEFT-SIDE HEADER"
r.font.size = Pt(16)
r.font.bold = True

add_bold_bullets(slide, [
    ("First left point", " supporting detail for left side"),
    ("Second left point", " more evidence for the left column"),
    ("Third left point", " additional data for this side"),
], left=0.67, top=2.1, width=5.5, height=4.0)

# Right header
rh = slide.shapes.add_textbox(Inches(6.5), Inches(1.6), Inches(5.5), Inches(0.4))
r = rh.text_frame.paragraphs[0].add_run()
r.text = "RIGHT-SIDE HEADER"
r.font.size = Pt(16)
r.font.bold = True

add_bold_bullets(slide, [
    ("First right point", " supporting detail for right side"),
    ("Second right point", " more evidence for the right column"),
    ("Third right point", " additional data for this side"),
], left=6.5, top=2.1, width=5.5, height=4.0)

add_source(slide)
print("two-column done")


# === 3. highlights ===
slide = prs.slides.add_slide(content_layout)
set_title(slide, "Sentence-form title showcasing key capabilities or achievements")
add_tag(slide, "Tag")

hdr = slide.shapes.add_textbox(Inches(0.67), Inches(1.6), Inches(5.0), Inches(0.4))
r = hdr.text_frame.paragraphs[0].add_run()
r.text = "HIGHLIGHTS"
r.font.size = Pt(16)
r.font.bold = True

for i, text in enumerate([
    "First key highlight statement with specific metric or proof point",
    "Second highlight statement with another independent achievement",
    "Third highlight statement with a differentiating feature and evidence",
]):
    add_callout_box(slide, text, 0.67, 2.2 + i * 1.3, 11.8, 1.0)

add_source(slide)
print("highlights done")


# === 4. data-table ===
slide = prs.slides.add_slide(content_layout)
set_title(slide, "Sentence-form title summarizing the data below")
add_tag(slide, "Tag")

rows, cols = 5, 4
tbl = slide.shapes.add_table(rows, cols, Inches(0.67), Inches(1.6), Inches(11.8), Inches(4.5)).table
for j, h in enumerate(["Column A", "Column B", "Column C", "Column D"]):
    tbl.cell(0, j).text = h
for i in range(1, rows):
    for j in range(cols):
        tbl.cell(i, j).text = f"Data {i},{j+1}"

add_source(slide)
print("data-table done")


# === 5. sidebar ===
slide = prs.slides.add_slide(content_layout)
set_title(slide, "Sentence-form title with main content and supporting callout")
add_tag(slide, "Tag")

tbl = slide.shapes.add_table(5, 3, Inches(0.67), Inches(1.6), Inches(7.5), Inches(4.5)).table
for j, h in enumerate(["Metric", "Year 1", "Year 2"]):
    tbl.cell(0, j).text = h
for i in range(1, 5):
    tbl.cell(i, 0).text = f"Row {i}"
    for j in range(1, 3):
        tbl.cell(i, j).text = "Data"

add_callout_box(slide, "Sidebar callout with pricing model, methodology note, or key metric summary", 8.5, 1.6, 3.9, 4.5)
add_source(slide)
print("sidebar done")


# === 6. testimonials ===
slide = prs.slides.add_slide(content_layout)
set_title(slide, "Sentence-form title highlighting the strongest proof point")
add_tag(slide, "Customer")

for i, (name, role) in enumerate([("Person Name", "Title, Company"), ("Second Person", "Title, Company")]):
    x = 0.67 + i * 6.0
    add_callout_box(slide, "Evidence point one with metric\nEvidence point two with data\nEvidence point three with proof", x, 1.6, 5.7, 3.5)
    attr = slide.shapes.add_textbox(Inches(x), Inches(5.3), Inches(5.7), Inches(0.5))
    r = attr.text_frame.paragraphs[0].add_run()
    r.text = f"{name}, {role}"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x00, 0xA0, 0xAC)

print("testimonials done")


# === 7. timeline ===
slide = prs.slides.add_slide(content_layout)
set_title(slide, "Sentence-form title describing the chronological progression")
add_tag(slide, "Timeline & Future Plan")

# Horizontal connector
slide.shapes.add_connector(1, Inches(1.0), Inches(2.7), Inches(11.5), Inches(2.7))

for i, year in enumerate(["2020", "2021", "2022", "2023", "2024"]):
    x = 0.8 + i * 2.4
    ym = slide.shapes.add_textbox(Inches(x), Inches(2.3), Inches(1.5), Inches(0.4))
    r = ym.text_frame.paragraphs[0].add_run()
    r.text = year
    r.font.size = Pt(18)
    r.font.bold = True

    ev = slide.shapes.add_textbox(Inches(x), Inches(3.0), Inches(2.0), Inches(1.5))
    ev.text_frame.word_wrap = True
    r = ev.text_frame.paragraphs[0].add_run()
    r.text = f"Milestone event for {year}"
    r.font.size = Pt(12)

print("timeline done")


# === 8. matrix ===
slide = prs.slides.add_slide(content_layout)
set_title(slide, "Sentence-form title positioning the company as strong on both axes")
add_tag(slide, "Competitive Analysis")

# Grid lines
slide.shapes.add_connector(1, Inches(6.3), Inches(1.8), Inches(6.3), Inches(6.2))
slide.shapes.add_connector(1, Inches(1.5), Inches(4.0), Inches(11.2), Inches(4.0))

# Axis labels
for text, x, y in [
    ("High Dimension A", 5.0, 1.4),
    ("Low Dimension A", 5.2, 6.3),
    ("Low Dimension B", 0.5, 3.8),
    ("High Dimension B", 10.5, 3.8),
]:
    lbl = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.5), Inches(0.4))
    r = lbl.text_frame.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(11)

# Quadrant entries
for text, x, y, bold in [
    ("COMPANY (winning)", 8.5, 2.2, True),
    ("Competitor A", 2.5, 2.5, False),
    ("Competitor B", 8.0, 4.5, False),
    ("Competitor C", 3.0, 4.8, False),
]:
    ent = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.5), Inches(0.8))
    ent.text_frame.word_wrap = True
    r = ent.text_frame.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.bold = bold

print("matrix done")


# === Step 3: Reorder and update S5/S6 ===
# After keeping S1-S4, S21, S22 and adding patterns, the order is:
#   S1, S2, S3, S4, [eBots S21], [eBots S22], [patterns...]
# We need to move S21 and S22 to the end so patterns are grouped after S4.

# Move slides using XML: sldIdLst controls slide order
sldIdLst = prs._element.find(f"{ns}sldIdLst")
slide_ids = list(sldIdLst)

# slide_ids[4] = eBots S21 (Investment Criteria)
# slide_ids[5] = eBots S22 (Appendix)
# Move both to the end
s21_elem = slide_ids[4]
s22_elem = slide_ids[5]
sldIdLst.remove(s21_elem)
sldIdLst.remove(s22_elem)
sldIdLst.append(s21_elem)
sldIdLst.append(s22_elem)

print("Reordered: S5 and S6 moved after pattern slides")

# Now S5 and S6 are the last two slides
num_slides = len(prs.slides)
s5 = prs.slides[num_slides - 2]
s6 = prs.slides[num_slides - 1]

# Update S5
for sh in s5.shapes:
    try:
        if sh.placeholder_format.idx == 0:
            replace_run_text(sh, "Investment criteria: Enter company name here")
    except (ValueError, AttributeError):
        pass
    if sh.has_table:
        table = sh.table
        for row in table.rows:
            if len(row.cells) >= 2:
                key = row.cells[0].text.strip()
                if key and key != "Attractiveness of Investment":
                    cell = row.cells[1]
                    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                        cell.text_frame.paragraphs[0].runs[0].text = f"Enter assessment for {key} here"
                        for run in cell.text_frame.paragraphs[0].runs[1:]:
                            run._r.getparent().remove(run._r)
                        for p in cell.text_frame.paragraphs[1:]:
                            p._p.getparent().remove(p._p)
print(f"S5 (Investment Criteria) done - layout={s5.slide_layout.name}, shapes={len(s5.shapes)}")

# S6 kept as-is
print(f"S6 (Appendix) done - layout={s6.slide_layout.name}, shapes={len(s6.shapes)}")

# === Final output ===
print(f"\nTotal: {len(prs.slides)} slides")
for i, sl in enumerate(prs.slides):
    title = ""
    for sh in sl.shapes:
        try:
            if sh.placeholder_format.idx == 0:
                title = sh.text[:50]
        except (ValueError, AttributeError):
            pass
    print(f"  {i+1:2d}. {sl.slide_layout.name:15s} | {title}")

prs.save(".claude/skills/investment-memo/assets/reference.pptx")
print("\nSaved reference.pptx")
