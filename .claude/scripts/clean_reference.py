"""
Clean reference.pptx so every slide is independent — no OLE objects,
no p:tags metadata, no orphan-prone relationships.  Only the logo image
on slide 1 is kept as a real relationship.

Run from repo root:
    python .claude/scripts/clean_reference.py
"""
import copy, os, sys
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

ASSET_DIR = os.path.join(os.path.dirname(__file__), '..', 'skills', 'investment-memo', 'assets')
REF_IN  = os.path.join(ASSET_DIR, 'reference.pptx')
REF_TMP = os.path.join(ASSET_DIR, 'reference_clean.pptx')

NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def strip_slide(slide, slide_num):
    """Remove OLE objects, tags, and custDataLst from a slide.
    Keep pictures (logo) and all normal shapes."""
    el = slide._element
    removed = []

    # 1. Remove graphicFrame elements that contain oleObj (think-cell)
    spTree = el.find(f'{{{NS_P}}}cSld').find(f'{{{NS_P}}}spTree')
    if spTree is None:
        # Try alternate path
        cSld = el.find(qn('p:cSld'))
        if cSld is not None:
            spTree = cSld.find(qn('p:spTree'))

    if spTree is not None:
        for gf in list(spTree.findall(qn('p:graphicFrame'))):
            # Check if it contains an oleObj
            has_ole = False
            for node in gf.iter():
                if node.tag == qn('p:oleObj') or 'oleObj' in node.tag:
                    has_ole = True
                    break
            if has_ole:
                spTree.remove(gf)
                removed.append('graphicFrame(OLE)')

    # 2. Remove all <p:tags> elements everywhere in the slide
    for tags_el in list(el.iter(qn('p:tags'))):
        parent = tags_el.getparent()
        if parent is not None:
            parent.remove(tags_el)
            removed.append('p:tags')

    # 3. Remove empty <p:custDataLst> containers
    for cdl in list(el.iter(qn('p:custDataLst'))):
        if len(cdl) == 0:
            parent = cdl.getparent()
            if parent is not None:
                parent.remove(cdl)
                removed.append('p:custDataLst(empty)')

    # 4. Drop orphaned relationships from this slide's rels
    #    After removing OLE/tags, some rIds are no longer referenced.
    #    Walk the slide XML, collect remaining rId refs, drop the rest.
    used_rids = set()
    for node in el.iter():
        for attr_name, attr_val in node.attrib.items():
            if attr_val.startswith('rId'):
                used_rids.add(attr_val)

    # Also collect rIds from <a:blip r:embed="...">
    for blip in el.iter(qn('a:blip')):
        embed = blip.get(qn('r:embed'))
        if embed:
            used_rids.add(embed)
        link = blip.get(qn('r:link'))
        if link:
            used_rids.add(link)

    # Drop unused rels
    for rel_key in list(slide.part.rels.keys()):
        if rel_key not in used_rids:
            rel = slide.part.rels[rel_key]
            rtype = rel.reltype.split('/')[-1]
            # Always keep slideLayout relationship
            if rtype == 'slideLayout':
                continue
            slide.part.drop_rel(rel_key)
            removed.append(f'rel:{rel_key}({rtype})')

    return removed


def main():
    prs = Presentation(REF_IN)
    print(f"Loaded: {REF_IN}")
    print(f"Slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides):
        removed = strip_slide(slide, i + 1)
        if removed:
            print(f"Slide {i+1:2d}: removed {removed}")
        else:
            print(f"Slide {i+1:2d}: already clean")

    # Verify: re-audit all slides
    print("\n--- Post-clean verification ---")
    all_clean = True
    for i, slide in enumerate(prs.slides):
        el = slide._element
        issues = []

        # Check for remaining OLE
        for node in el.iter():
            if 'oleObj' in node.tag:
                issues.append('oleObj')
                break

        # Check for remaining tags
        for node in el.iter(qn('p:tags')):
            issues.append('p:tags')
            break

        # Check rels
        for rel_key in slide.part.rels:
            rel = slide.part.rels[rel_key]
            rtype = rel.reltype.split('/')[-1]
            if rtype not in ('slideLayout', 'image', 'notesSlide', 'hyperlink'):
                issues.append(f'rel:{rel_key}({rtype})')

        if issues:
            print(f"  Slide {i+1}: ISSUES {issues}")
            all_clean = False
        else:
            # Count remaining non-layout rels
            non_layout = [k for k in slide.part.rels if slide.part.rels[k].reltype.split('/')[-1] != 'slideLayout']
            suffix = f" ({len(non_layout)} non-layout rels)" if non_layout else ""
            print(f"  Slide {i+1}: CLEAN{suffix}")

    if all_clean:
        prs.save(REF_TMP)
        print(f"\nSaved clean reference to: {REF_TMP}")
        print(f"Replace original: copy {REF_TMP} -> {REF_IN}")
    else:
        prs.save(REF_TMP)
        print(f"\nSaved with warnings to: {REF_TMP}")


if __name__ == '__main__':
    main()
